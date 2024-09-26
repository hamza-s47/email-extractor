import PyPDF2
import textract
from docx import Document
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from openpyxl import load_workbook
from io import BytesIO

# URL Extractor
def extract_url(url):
    try:
        response = requests.get(url)
        
        if response.status_code != 200:
            raise Exception(f"Failed to retrieve the page: {response.status_code}")
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        return soup.get_text(separator=' ')
    
    except Exception as e:
        return f"Error extracting text from URL: {e}"

# PDF Extractor
def extract_pdf(data):
        reader = PyPDF2.PdfReader(data.stream)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
        return text

# DOCX Extractor
def extract_docx(file):
    doc = Document(file)
    return '\n'.join([para.text for para in doc.paragraphs])

# DOC Extractor(using docx2txt)
def extract_doc(file):
    text = textract.process(file, extension='doc').decode('utf-8')
    return text

# PPTX Extractor
def extract_pptx(file):
    pptx_file = BytesIO(file.read())
    prs = Presentation(pptx_file)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

# TXT Extractor
def extract_txt(file):
    text = file.stream.read().decode('utf-8')
    return text

# XLSX Extractor
def extract_excel(file):
    xlsx_file = BytesIO(file.read())
    workbook = load_workbook(xlsx_file, data_only=True)
    text = ''
    
    for sheet in workbook.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                text += row_text + '\n'

    return text

# Master Function
def extract_text_from_file(file):
    
    if file.filename.endswith('.pdf'):
        return extract_pdf(file)
    elif file.filename.endswith('.docx'):
        return extract_docx(file)
    elif file.filename.endswith('.doc'):
        return extract_doc(file)
    elif file.filename.endswith('.pptx'):
        return extract_pptx(file)
    elif file.filename.endswith('.txt'):
        return extract_txt(file)
    elif file.filename.endswith('.xlsx'):
        return extract_excel(file)
    else:
        raise ValueError(f"Unsupported file type: {file.filename}")